import glob
from pptx import Presentation
from pptx.util import Inches
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]
for f in sorted(glob.glob('slides_png/slide-*.png')):
    s = prs.slides.add_slide(blank)
    s.shapes.add_picture(f, 0, 0, width=prs.slide_width, height=prs.slide_height)
prs.save('Tax-AI-Keynote-Google.pptx')
print('saved Tax-AI-Keynote-Google.pptx with', len(prs.slides._sldIdLst), 'slides')
